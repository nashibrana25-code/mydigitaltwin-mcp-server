import os
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches

# Create a new presentation
prs = Presentation()

# Set slide size to 16:9 aspect ratio (10 inches wide, 5.625 inches high)
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

with sync_playwright() as p:
    browser = p.chromium.launch()
    for i in range(1, 17):
        page = browser.new_page()
        page.set_viewport_size({"width": 1280, "height": 720})
        html_file = f"c:/Users/nashi/Week 5/slide{i}.html"
        page.goto(f"file://{html_file}")
        page.wait_for_load_state('networkidle')
        screenshot_path = f"temp_slide_{i}.png"
        page.screenshot(path=screenshot_path, full_page=True)
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank slide layout
        left = top = Inches(0)
        width = Inches(10)
        height = Inches(5.625)
        slide.shapes.add_picture(screenshot_path, left, top, width, height)
        os.remove(screenshot_path)  # Clean up temporary file
    browser.close()

# Save the presentation
prs.save("presentation.pptx")